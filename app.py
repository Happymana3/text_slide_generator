import streamlit as st
from transformers import pipeline
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import PyPDF2
import pandas as pd

# ✅ Set page config FIRST
st.set_page_config(page_title="Slide Generator", layout="centered")

st.title("📊 AI Slide Generator")
st.write("Upload a `.txt`, `.pdf`, or `.xlsx` file to generate presentation slides.")

# Load summarizer
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

def summarize_text(text, max_len=200):
    chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
    summaries = [summarizer(chunk, max_length=max_len, min_length=30, do_sample=False)[0]["summary_text"] for chunk in chunks]
    return summaries

def create_ppt(summaries):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[1]

    for i, summary in enumerate(summaries):
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Set title
        title.text = f"Slide {i + 1}"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

        # Format bullet points
        content.text = ""  # Clear default text
        for bullet in summary.split('. '):
            if bullet.strip():  # Avoid empty strings
                p = content.text_frame.add_paragraph()
                p.text = bullet.strip().rstrip('.') + '.'
                p.level = 0
                p.font.size = Pt(20)
                p.font.name = 'Calibri'
                p.font.color.rgb = RGBColor(50, 50, 50)  # Dark gray

    return prs

uploaded_file = st.file_uploader("Upload File", type=["txt", "pdf", "xlsx"])

if uploaded_file is not None:
    if uploaded_file.name.endswith(".txt"):
        text = uploaded_file.read().decode("utf-8")

    elif uploaded_file.name.endswith(".pdf"):
        reader = PyPDF2.PdfReader(uploaded_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text()

    elif uploaded_file.name.endswith(".xlsx"):
        df = pd.read_excel(uploaded_file)
        text = df.astype(str).apply(lambda x: ' '.join(x), axis=1).str.cat(sep=' ')

    else:
        st.error("Unsupported file type.")
        text = ""

    if text:
        st.success("✅ File processed. Generating slides...")
        summaries = summarize_text(text)
        ppt = create_ppt(summaries)

        ppt_io = BytesIO()
        ppt.save(ppt_io)
        ppt_io.seek(0)

        st.success("✅ Slides generated!")
        st.download_button(
            label="📥 Download Slides as .pptx",
            data=ppt_io,
            file_name="generated_slides.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
